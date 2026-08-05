from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

path = sys.argv[1]
layout_ids = sys.argv[2].split(',') if len(sys.argv) > 2 else []
prs = Presentation(path)
print(f'Slides: {len(prs.slides)}')
print(f'Slide width: {prs.slide_width.inches:.3f} in, height: {prs.slide_height.inches:.3f} in')
blank_slides = []
off_slides = []
for i, slide in enumerate(prs.slides, 1):
    shapes = list(slide.shapes)
    texts = []
    charts = []
    images = []
    off = []
    for s in shapes:
        txt = ''
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                txt += p.text
        if txt.strip():
            texts.append(txt[:60])
        if s.shape_type == MSO_SHAPE_TYPE.CHART:
            try:
                charts.append(f'{s.chart.chart_type} series={len(s.chart.series)}')
            except Exception as e:
                charts.append(f'chart (error: {e})')
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append('image')
        x = s.left.inches if s.left else 0
        y = s.top.inches if s.top else 0
        w = s.width.inches if s.width else 0
        h = s.height.inches if s.height else 0
        sw = prs.slide_width.inches
        sh = prs.slide_height.inches
        if x+w > sw+0.05 or y+h > sh+0.05 or x < -0.05 or y < -0.05:
            off.append(f'{s.shape_type} at {x:.2f},{y:.2f} {w:.2f}x{h:.2f}')
    layout = layout_ids[i-1] if i-1 < len(layout_ids) else '?'
    is_blank = len(texts) == 0 and len(charts) == 0 and len(images) == 0
    if is_blank:
        blank_slides.append((i, layout))
    if off:
        off_slides.append((i, layout, off))
    print(f'\nSlide {i} ({layout}): shapes={len(shapes)}, text={len(texts)}, charts={len(charts)}, images={len(images)}')
    if is_blank:
        print('  [BLANK]')
    else:
        for t in texts[:6]:
            print(f'  text: {t}')
        for c in charts:
            print(f'  chart: {c}')
        for img in images[:3]:
            print(f'  image: {img}')
    if off:
        print('  OFF-SLIDE:')
        for o in off[:8]:
            print(f'    {o}')

print('\n--- SUMMARY ---')
if blank_slides:
    print(f'Blank slides ({len(blank_slides)}):')
    for i, layout in blank_slides:
        print(f'  {i}: {layout}')
else:
    print('No blank slides.')
if off_slides:
    print(f'Slides with off-slide shapes ({len(off_slides)}):')
    for i, layout, off in off_slides:
        print(f'  {i}: {layout} ({len(off)} shapes)')
else:
    print('No off-slide shapes.')
